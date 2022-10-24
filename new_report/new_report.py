import datetime

import dateutil
import docx
import traceback
import asyncio

import httpx
from docx.oxml.shared import qn
from docx.oxml import OxmlElement
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, ChartData
from docx import Document
from docx.enum.style import WD_STYLE_TYPE
from pptx.util import Pt, Inches
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from docx.shared import Mm

from new_report.request_file import getSourceStats, get_additional_info, get_trustdaily, get_stats, get_ages, get_city, \
    get_top, get_post_top
from word_media import login

COOKIES = []

KOM_NAME = "Комитет по образованию"
STYLE = "Times New Roman"
PT = Pt(10.5)


def update_chart_style(chart):
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    chart.legend.font.name = STYLE

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(10)
    data_labels.font.name = STYLE
    data_labels._element.get_or_add_txPr().bodyPr.set('rot', '-5400000')

    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = RGBColor(0, 0, 0)
    value_axis.format.line.color.rgb = RGBColor(0, 0, 0)
    chart.category_axis.format.line.color.rgb = RGBColor(0, 0, 0)
    value_axis.format.line.color.rgb = RGBColor(0, 0, 0)
    value_axis.format.line.fill.fore_color.rgb = RGBColor(0, 0, 0)
    value_axis.major_gridlines.format.line.fill.fore_color.rgb = RGBColor(211, 211, 211)
    value_axis.major_gridlines.format.line.color.rgb = RGBColor(211, 211, 211)
    value_axis.tick_labels.font.color.rgb = RGBColor(211, 211, 211)
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.name = STYLE
    chart.value_axis.tick_labels.font.size = Pt(1)
    chart.value_axis.tick_labels.font.fill.rgb = RGBColor(150, 150, 150)
    chart.value_axis.tick_labels.font.name = STYLE
    chart.plots[0].chart.value_axis.tick_labels.font.size = Pt(1)
    chart.plots[0].chart.value_axis.tick_labels.font.name = STYLE
    chart.plots[0].chart.category_axis.tick_labels.font.size = Pt(10)
    chart.plots[0].chart.category_axis.tick_labels.font.name = STYLE

    plot.overlap = -10


def add_pie_chart(document, data):
    chart_data = ChartData()
    chart_data.categories = ['СМИ', 'Вконтакте', 'FaceBook', 'Twitter', 'Telegram', "Youtube"]
    chart_data.add_series('Series 1',
                          get_list_([data["gs"]["total"],
                                     data["vk"]["total"],
                                     data["fb"]["total"],
                                     data["tw"]["total"],
                                     data["tg"]["total"],
                                     data["yt"]["total"]])
                          )
    x, y, cx, cy = Inches(-1.0), Inches(0), Inches(6.55), Inches(3.3)
    chart = document.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data)

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.name = STYLE
    chart.legend.font.size = Pt(10.5)

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0.0%'
    data_labels.font.name = STYLE
    data_labels.font.size = Pt(10.5)

    color_list = [
        RGBColor(217, 150, 148),
        RGBColor(147, 205, 221),
        RGBColor(49, 133, 156),
        RGBColor(204, 193, 218),
        RGBColor(196, 189, 151),
        RGBColor(250, 192, 144)

    ]
    for idx, point in enumerate(chart.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color_list[idx]
    p = document.paragraphs[-1].add_run("Рисунок 1. Распределение публикаций по предоставленным источникам",
                                        style=STYLE)
    p.italic = True
    document.paragraphs[-1].alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER


def add_double_chart(document, dates, smi, social):
    import locale
    locale.setlocale(locale.LC_TIME, 'ru_RU.UTF-8')

    date_str = []
    for d in dates:
        d_date_time = datetime.datetime.strptime(d, "%Y-%m-%d")
        date_str.append(d_date_time.strftime('%d-%b'))

    chart_data = CategoryChartData()
    chart_data.categories = date_str
    chart_data.add_series('Сми', list(smi.values()))
    chart_data.add_series('Социальные сети', list(social.values()))
    x, y, cx, cy = Inches(-1.0), Inches(0), Inches(6.55), Inches(3.3)

    chart = document.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
    change_color(chart.plots[0].series[0], RGBColor(228, 108, 10))
    change_color(chart.plots[0].series[1], RGBColor(23, 55, 94))

    update_chart_style(chart)
    p = document.paragraphs[-1].add_run("Рисунок 2. Динамика публикаций по предоставленным источникам",
                                        style=STYLE)
    p.italic = True
    document.paragraphs[-1].alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER


def add_double_vk_tg_chart(document):
    chart_data = CategoryChartData()
    chart_data.categories = ['9/11/2022', '10/11/2022', '11/11/2022', '12/11/2022']
    chart_data.add_series('ВКонтакте', (6202108, 7802108, 10902108, 8099999))
    chart_data.add_series('Telegram', (26202108, 89202108, 99202108, 59202108))
    x, y, cx, cy = Inches(-3.5), Inches(0), Inches(6.15), Inches(3.3)

    chart = document.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
    change_color(chart.plots[0].series[0], RGBColor(255, 34, 13))
    change_color(chart.plots[0].series[1], RGBColor(0, 32, 255))

    update_chart_style(chart)


def add_triple_hart(document, cat, like_, repost_, comment_):
    date_str = []
    for d in cat:
        d_date_time = datetime.datetime.strptime(d, "%Y-%m-%d")
        date_str.append(d_date_time.strftime('%d-%b'))
    chart_data = CategoryChartData()
    chart_data.categories = date_str
    chart_data.add_series('Лайки', like_)
    chart_data.add_series('Репосты', repost_)
    chart_data.add_series('Комментарии', comment_)

    x, y, cx, cy = Inches(-1.0), Inches(0), Inches(6.55), Inches(3.3)

    chart = document.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
    change_color(chart.plots[0].series[0], RGBColor(217, 150, 148))
    change_color(chart.plots[0].series[1], RGBColor(85, 142, 213))
    change_color(chart.plots[0].series[2], RGBColor(55, 96, 146))

    update_chart_style(chart)
    p = document.paragraphs[-1].add_run(
        "Рисунок 4. Динамика активности лайков, репостов и комментариев в социальных сетях",
        style=STYLE)
    p.italic = True
    document.paragraphs[-1].alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER


def get_list_(_list):
    all = sum(_list)
    res = []
    if all ==0:
        return _list
    for l in _list:
        res.append(f"{l / all:.{3}f}")
    return res


def add_pie_chart_stat(document, netural, negative, positive):
    chart_data = ChartData()

    chart_data.categories = ['Негативные публикации', 'Нейтральные публикации', 'Позитивные публикации']
    chart_data.add_series('Series 1', get_list_([negative, netural, positive]))
    x, y, cx, cy = Inches(-1.0), Inches(0), Inches(6.55), Inches(3.3)
    chart = document.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data)

    # chart = slide.shapes.add_chart(
    #     XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    # ).chart
    #
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.name = STYLE
    chart.legend.font.size = Pt(9)

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    # data_labels.number_format = '0.0%'
    data_labels.number_format_is_linked = False
    data_labels.number_format = '0%'
    data_labels.font.name = STYLE
    data_labels.font.size = Pt(10.5)

    color_list = [
        # RGBColor.from_string(color_list[col_idx])
        RGBColor(217, 150, 148),
        RGBColor(149, 179, 215),
        RGBColor(195, 214, 155),
    ]
    for idx, point in enumerate(chart.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color_list[idx]
    p = document.paragraphs[-1].add_run("Рисунок 5. Тональность публикаций",
                                        style=STYLE)
    p.italic = True
    document.paragraphs[-1].alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER


def add_pie_age(document, group_1, group_2, group_3, group_4):
    chart_data = ChartData()

    chart_data.categories = ['18-25 лет', '26-39 лет', '40-54 лет', '55+ лет']
    chart_data.add_series('Series 1', get_list_([group_1, group_2, group_3, group_4]))
    x, y, cx, cy = Inches(-1.0), Inches(0), Inches(6.55), Inches(3.3)
    chart = document.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data)

    # chart = slide.shapes.add_chart(
    #     XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    # ).chart
    #
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.name = STYLE
    chart.legend.font.size = Pt(9)

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0.0%'
    data_labels.font.name = STYLE
    data_labels.font.size = Pt(10.5)

    color_list = [
        # RGBColor.from_string(color_list[col_idx])
        RGBColor(149, 179, 215),
        RGBColor(252, 213, 181),
        RGBColor(195, 214, 155),
        RGBColor(85, 142, 213),

    ]
    for idx, point in enumerate(chart.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color_list[idx]
    p = document.paragraphs[-1].add_run("Рисунок 6. Возраст активной аудитории",
                                        style=STYLE)
    p.italic = True
    document.paragraphs[-1].alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER


def add_pie_sex(document, m, w):
    chart_data = ChartData()

    chart_data.categories = ['Мужской', 'Женский']
    chart_data.add_series('Series 1', get_list_([m, w]))
    x, y, cx, cy = Inches(-1.0), Inches(0), Inches(6.55), Inches(3.3)
    chart = document.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data)

    # chart = slide.shapes.add_chart(
    #     XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    # ).chart
    #
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.name = STYLE
    chart.legend.font.size = Pt(9)

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0.0%'
    data_labels.font.name = STYLE
    data_labels.font.size = Pt(10.5)

    color_list = [
        # RGBColor.from_string(color_list[col_idx])
        RGBColor(149, 179, 215),
        RGBColor(230, 185, 184),

    ]
    for idx, point in enumerate(chart.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color_list[idx]
    p = document.paragraphs[-1].add_run("Рисунок 7. Пол активной аудитории",
                                        style=STYLE)
    p.italic = True
    document.paragraphs[-1].alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER


def add_pie_city(document, spb, anpther):
    chart_data = ChartData()

    chart_data.categories = ['Санкт-Петербург', 'Подписчики из другого города']
    chart_data.add_series('Series 1', get_list_([spb, anpther]))
    x, y, cx, cy = Inches(-1.0), Inches(0), Inches(6.55), Inches(3.3)
    chart = document.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data)

    # chart = slide.shapes.add_chart(
    #     XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
    # ).chart
    #
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.name = STYLE
    chart.legend.font.size = Pt(9)

    chart.plots[0].has_data_labels = True
    data_labels = chart.plots[0].data_labels
    data_labels.number_format = '0.0%'
    data_labels.font.name = STYLE
    data_labels.font.size = Pt(10.5)

    color_list = [
        # RGBColor.from_string(color_list[col_idx])
        RGBColor(149, 179, 215),
        RGBColor(195, 214, 155),

    ]
    for idx, point in enumerate(chart.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = color_list[idx]
    p = document.paragraphs[-1].add_run("Рисунок 8. Геолокация активной аудитории",
                                        style=STYLE)
    p.italic = True
    document.paragraphs[-1].alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER


def change_color(series, color):
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_table(document, title, last_colum, items):
    table = document.add_table(rows=1, cols=4, style='Table Grid')

    table.columns[0].width = Inches(0.5)
    table.columns[1].width = Inches(2.4)
    table.columns[2].width = Inches(2.4)
    table.columns[3].width = Inches(1.3)
    hdr_cells = table.rows[0].cells

    hdr_cells[0].text = title
    hdr_cells[0].style = STYLE
    hdr_cells[0].paragraphs[0].runs[0].bold = True

    hdr_cells[0].paragraphs[0].alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER

    table.rows[0].cells[0].merge(table.rows[0].cells[-1])
    for i in range(4):
        set_cell_vertical_alignment(hdr_cells[i], align="center")
        # hdr_cells[i].paragraphs[-1].runs[-1].style = STYLE
        # hdr_cells[i].paragraphs[-1].runs[-1].size = docx.shared.Pt(12)
    hdr_cells[0].paragraphs[-1].runs[-1].style = STYLE
    hdr_cells[0].paragraphs[-1].runs[-1].size = docx.shared.Pt(12)

    row_cells = table.add_row().cells
    row_cells[0].text = "№ п/п"
    row_cells[1].text = "Источник"
    row_cells[2].text = "Ссылка"
    row_cells[3].text = last_colum
    for i in range(4):
        # row_cells[i].paragraphs[0].style.name = STYLE
        # row_cells[i].paragraphs[0].style.name = STYLE
        row_cells[i].paragraphs[-1].runs[-1].style = STYLE
        row_cells[i].paragraphs[-1].runs[-1].size = docx.shared.Pt(12)

        row_cells[i].style = STYLE
        row_cells[i].paragraphs[0].runs[0].bold = True
        row_cells[i].paragraphs[0].alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER
        set_cell_vertical_alignment(row_cells[i], align="center")

    t = 1
    for item in items:
        row_cells = table.add_row().cells
        row_cells[0].text = str(t)
        row_cells[1].text = item['title']

        add_hyperlink(row_cells[2].paragraphs[0], item['url'], item['url'], None, True)

        row_cells[3].text = str(item['post_count'])
        for i in range(4):
            if i == 0 or i == 3:
                row_cells[i].paragraphs[0].alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER
            try:
                row_cells[i].paragraphs[-1].runs[-1].style = STYLE
                row_cells[i].paragraphs[-1].runs[-1].size = docx.shared.Pt(12)
            except Exception:
                pass
            row_cells[i].style = STYLE
            set_cell_vertical_alignment(row_cells[i], align="center")
        t += 1


def add_table_pril(document, items):
    table = document.add_table(rows=1, cols=6, style='Table Grid')

    table.columns[0].width = Inches(0.5)
    table.columns[1].width = Inches(0.7)
    table.columns[2].width = Inches(1.9)
    table.columns[3].width = Inches(1.4)
    table.columns[4].width = Inches(0.9)
    table.columns[5].width = Inches(1.2)

    row_cells = table.rows[0].cells
    row_cells[0].text = "№ п/п"
    row_cells[1].text = "Соц сеть"
    row_cells[2].text = "Пользователь/группа (разместивший пост)"
    row_cells[3].text = "Ссылка на пост"
    row_cells[4].text = "Дата"
    row_cells[5].text = "Охват (пользователи)"

    for i in range(6):
        row_cells[i].style = STYLE
        row_cells[i].paragraphs[0].runs[0].bold = True
        row_cells[i].paragraphs[0].alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER
        set_cell_vertical_alignment(row_cells[i], align="center")
        row_cells[i].paragraphs[-1].runs[-1].style = STYLE
        row_cells[i].paragraphs[-1].runs[-1].size = docx.shared.Pt(12)
    row_cells = table.add_row().cells
    row_cells[0].text = "1"
    row_cells[1].text = "2"
    row_cells[2].text = "3"
    row_cells[3].text = "4"
    row_cells[4].text = "5"
    row_cells[5].text = "6"

    for i in range(6):
        row_cells[i].style = STYLE
        row_cells[i].paragraphs[0].alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER
        set_cell_vertical_alignment(row_cells[i], align="center")
        row_cells[i].paragraphs[-1].runs[-1].style = STYLE
        row_cells[i].paragraphs[-1].runs[-1].size = docx.shared.Pt(12)
    t = 1
    for item in items:
        row_cells = table.add_row().cells
        row_cells[0].text = str(t)
        if item['network_name'] == "gs":
            row_cells[1].text = "СМИ"
        else:
            row_cells[1].text = item['network_name']
        row_cells[2].text = item['author']
        add_hyperlink(row_cells[3].paragraphs[0], item['url'], item['url'], None, True)
        row_cells[4].text = get_date(item['created_date']).strftime("%d.%m.%Y")
        row_cells[5].text = get_str_int(item['attendance'])
        for i in range(6):
            row_cells[i].style = STYLE
            try:
                row_cells[i].paragraphs[-1].runs[-1].style = STYLE
                row_cells[i].paragraphs[-1].runs[-1].size = docx.shared.Pt(12)
            except Exception:
                pass
            if i == 0 or i == 1 or i == 5:
                row_cells[i].paragraphs[0].alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER
            set_cell_vertical_alignment(row_cells[i], align="center")
        t += 1


def set_cell_vertical_alignment(cell, align="center"):
    try:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        tcValign = OxmlElement('w:vAlign')
        tcValign.set(qn('w:val'), align)
        tcPr.append(tcValign)
        return True
    except Exception:
        traceback.print_exc()
        return False


DEFAULT_DATE_FORMAT = "%Y-%m-%d %H:%M:%S"


def add_hyperlink(paragraph, url, text, color, underline, is_italic=False):
    part = paragraph.part
    r_id = part.relate_to(url, docx.opc.constants.RELATIONSHIP_TYPE.HYPERLINK, is_external=True)

    hyperlink = docx.oxml.shared.OxmlElement('w:hyperlink')
    hyperlink.set(docx.oxml.shared.qn('r:id'), r_id, )

    new_run = docx.oxml.shared.OxmlElement('w:r')

    rPr = docx.oxml.shared.OxmlElement('w:rPr')
    if is_italic:
        u = docx.oxml.shared.OxmlElement('w:i')
        u.append(docx.oxml.shared.OxmlElement('w:iCs'))
        rPr.append(u)
    if not color is None:
        c = docx.oxml.shared.OxmlElement('w:color')
        c.set(docx.oxml.shared.qn('w:val'), color)
        rPr.append(c)

    if not underline:
        u = docx.oxml.shared.OxmlElement('w:u')
        u.set(docx.oxml.shared.qn('w:val'), 'none')
        rPr.append(u)

    new_run.append(rPr)
    new_run.text = text
    new_run.style = STYLE
    hyperlink.append(new_run)
    paragraph._p.append(hyperlink)
    # paragraph._p.style = STYLE
    paragraph.style.font.name = "Times New Roman"
    paragraph.style.font.size = docx.shared.Pt(12)

    return hyperlink


def get_date(date):
    try:
        return datetime.datetime.strptime(date, DEFAULT_DATE_FORMAT)
    except Exception:
        try:
            return datetime.datetime.strptime(date, "%Y-%m-%d")
        except Exception:
            return dateutil.parser.parse(date)


DATES = {
    1: "января",
    2: "февраля",
    3: "марта",
    4: "апреля",
    5: "мая",
    6: "июня",
    7: "июля",
    8: "авуста",
    9: "сентября",
    10: "октября",
    11: "ноября",
    12: "декабря",

}


def get_str_int(x):
    return '{0:,}'.format(x).replace(',', ' ')


async def prepare_report(thread_id, _from, _to, _login, _password):
    async with httpx.AsyncClient() as session:
        await login(session)
        try:
            top_post, additional_info, json_stats, (count, res_dict, trust_daily), likes_stats, comments_stats, \
            reposts_stats, ages, city, top_smi, top_social = await asyncio.gather(
                get_post_top(session, thread_id, _from, _to),
                get_additional_info(session, thread_id),
                getSourceStats(session, thread_id, _from, _to),
                get_trustdaily(session, thread_id, _from, _to),
                get_stats(session, thread_id, _from, _to, "likes"),
                get_stats(session, thread_id, _from, _to, "comments"),
                get_stats(session, thread_id, _from, _to, "reposts"),
                get_ages(session, thread_id, _from, _to),
                get_city(session, thread_id, _from, _to),
                get_top(session, thread_id, _from, _to, "smi"),
                get_top(session, thread_id, _from, _to, "social")
            )
        except Exception as e:
            print(e)
            raise e
    document = Document("new_report/test.docx")

    obj_styles = document.styles
    obj_charstyle = obj_styles.add_style(STYLE, WD_STYLE_TYPE.CHARACTER)
    obj_font = obj_charstyle.font
    obj_font.size = Pt(12)
    obj_font.name = STYLE
    from_date = get_date(_from)
    to_date = get_date(_to)

    add_title_text(document,
                   f"Результаты сбора информации по направлению мониторинга «{additional_info.get('name')}» в "
                   f"онлайн-СМИ и социальных сетях за период с {from_date.day}  по {to_date.day} {DATES[to_date.month]} {to_date.year} года ",
                   True)

    new_paragraph(document, "Общее количество публикаций по каждой социальной сети", is_italic=True, paragraph=1)
    p = new_paragraph(document, paragraph=2)
    document.paragraphs[2].runs[
        0].text = f'За период с {from_date.day} {DATES[from_date.month]} {from_date.year} г. {to_date.day} {DATES[to_date.month]} {to_date.year} г. в СМИ и социальных сетях ВКонтакте,'
    p.add_run(', Twitter, Telegram и YouTube было размещено', style=STYLE)

    bold_text = p.add_run(f' {count} публикации.', style=STYLE)
    bold_text.bold = True

    new_paragraph(document, "Количество реакций пользователей на публикации", is_italic=True)

    p = new_paragraph(document)
    p.add_run('Все публикации в сумме набрали', style=STYLE)

    viewed = 0
    likes = 0
    reposts = 0
    comments = 0
    for stat in json_stats:
        viewed += stat.get("viewed")
        likes += stat.get("likes")
        reposts += stat.get("reposts")
        comments += stat.get("comments")

    bold_text = p.add_run(
        f' {get_str_int(viewed)} просмотр, {get_str_int(likes)} лайков, {get_str_int(reposts)} репостов и {get_str_int(comments)} комментария.',
        style=STYLE)
    bold_text.bold = True

    new_paragraph(document, "Распределение публикаций по предоставленным источникам", is_italic=True)

    p = new_paragraph(document)
    p.add_run('В СМИ была размещена ', style=STYLE)
    bold_text = p.add_run(f' {get_str_int(res_dict["gs"]["total"])} ', style=STYLE)
    bold_text.bold = True
    p.add_run(f' публикация. В социальной сети ВКонтакте – ', style=STYLE)
    bold_text = p.add_run(f' {get_str_int(res_dict["vk"]["total"])}', style=STYLE)
    bold_text.bold = True
    p.add_run(f' , в Facebook – ', style=STYLE)
    bold_text = p.add_run(f' {get_str_int(res_dict["fb"]["total"])}', style=STYLE)
    bold_text.bold = True
    p.add_run(f' ,в Twitter – ', style=STYLE)
    bold_text = p.add_run(f' {get_str_int(res_dict["tw"]["total"])}', style=STYLE)
    bold_text.bold = True
    p.add_run(f', в Telegram – ', style=STYLE)
    bold_text = p.add_run(f' {get_str_int(res_dict["tg"]["total"])}', style=STYLE)
    bold_text.bold = True
    p.add_run(f', в YouTube – ', style=STYLE)
    bold_text = p.add_run(f' {get_str_int(res_dict["yt"]["total"])}', style=STYLE)
    bold_text.bold = True
    p.add_run(f'.', style=STYLE)

    add_pie_chart(document, res_dict)

    new_paragraph(document,
                  "По исследуемой теме в СМИ и социальных сетях размещались материалы, посвященные ")
    new_paragraph(document,
                  "На рисунке 2 указана динамика публикаций за отчетный период, на рисунке 3 – динамика просмотров, на рисунке 4 – активность лайков, репостов и комментариев, на рисунке 5 – тональность публикаций, на рисунках 6 и 7 указано сегментирование активной аудитории по половозрастным категориям, на рисунке 8 – географическое расположение активной аудитории. ",
                  is_JUSTIFY=True)

    dates = set()
    smi = {}
    social = {}
    for tr in trust_daily:
        if tr != "total":
            for s in trust_daily.get(tr):
                if s != "total":
                    for d in trust_daily.get(tr).get(s):
                        dates.add(d[0])
                        if tr == "gs":
                            if social.get(d[0]):
                                social[d[0]] += d[-1]
                            else:
                                social[d[0]] = d[-1]

                        else:
                            if smi.get(d[0]):
                                smi[d[0]] += d[-1]
                            else:
                                smi[d[0]] = d[-1]
    dates = sorted(dates)
    add_double_chart(document, dates, smi, social)

    max_date_p = list(dates)[0]
    max_t = smi[max_date_p] + social[max_date_p]
    for d in list(dates)[1:]:
        if smi[d] + social[d] > max_t:
            max_t = smi[d] + social[d]
            max_date_p = d
    p = new_paragraph(document,
                      f"""Пиковый день активности по количеству публикаций приходится на """
                      )
    datetime_usr = datetime.datetime.strptime(max_date_p, "%Y-%m-%d")
    bold_text = p.add_run(f'{datetime_usr.day} {DATES[datetime_usr.month]} {datetime_usr.year} г.', style=STYLE)
    bold_text.bold = True
    p.add_run(
        f' "в СМИ – {get_str_int(smi[max_date_p])} публикаций, в социальных сетях – {get_str_int(social[max_date_p])} публикаций"',
        style=STYLE)
    new_paragraph(document, "Перечень основных информационных поводов публикаций", is_italic=True)

    reaction = {}
    for t in json_stats:
        d = t['viewed'] + t['likes'] + t['reposts'] + t['comments']
        if reaction.get(t['network_name']) is None:
            reaction[t['network_name']] = d
        else:
            reaction[t['network_name']] += d

    # add_double_vk_tg_chart(document)
    p = new_paragraph(document,
                      "Рисунок 3. Динамика активности просмотров в социальных сетях"
                      )
    p.italic = True
    p.alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER

    p = new_paragraph(document,
                      f"""Пиковый день активности просмотров приходится на  """
                      )
    p = new_paragraph(document,
                      f"""Перечень основных информационных поводов публикаций   """
                      )

    cat = []
    like_ = []
    repost_ = []
    comment_ = []
    for i in range(len(likes_stats['graph_data'])):
        cat.append(likes_stats['graph_data'][i][0])
        like_.append(likes_stats['graph_data'][i][-1])
        repost_.append(reposts_stats['graph_data'][i][-1])
        comment_.append(comments_stats['graph_data'][i][-1])

    add_triple_hart(document, cat, like_, repost_, comment_)

    max_date = 0
    max_t = like_[0] + repost_[0] + comment_[0]
    k = 0
    for d in range(len(cat)):
        if like_[d] + repost_[d] + comment_[d] > max_t:
            max_t = like_[d] + repost_[d] + comment_[d]
            max_date = d
    p = new_paragraph(document,
                      f"""Пиковый день активности лайков, репостов и комментариев приходится на  –"""
                      )
    datetime_usr = datetime.datetime.strptime(cat[max_date], "%Y-%m-%d")
    bold_text = p.add_run(f'{datetime_usr.day} {DATES[datetime_usr.month]} {datetime_usr.year} г. – ', style=STYLE)
    bold_text.bold = True
    p.add_run(
        f" {get_str_int(like_[max_date])} лайка, {get_str_int(repost_[max_date])} репоста,{get_str_int(comment_[max_date])} комментария. ",
        style=STYLE)

    new_paragraph(document, "Перечень основных информационных поводов публикаций", is_italic=True)
    new_paragraph(document, "Статистика эмоционального окраса в публикациях", is_italic=True)

    netural = 0
    negative = 0
    positive = 0
    for r in res_dict:
        netural += res_dict[r]['netural']
        negative += res_dict[r]['negative']
        positive += res_dict[r]['positive']

    add_pie_chart_stat(document, netural, negative, positive)

    new_paragraph(document, "ТОП-10 публикаций за сутки представлен в Приложении 1.", is_italic=True)

    new_paragraph(document, "ТОП-10 источников СМИ с наибольшим количеством публикуемой информации", is_italic=True)
    add_table(document, title="Топ-10 по количеству публикаций в СМИ", last_colum="Количество публикаций",
              items=top_smi['items'])

    new_paragraph(document, "ТОП-10 источников социальных сетей с наибольшим количеством публикуемой информации",
                  is_italic=True)
    add_table(document, title="Топ-10 по количеству публикаций в социальных сетях", last_colum="Количество публикаций",
              items=top_social['items'])

    new_paragraph(document,
                  "ТОП-10 публикаций по сумме просмотров в социальных сетях (оставлять просмотры можно только во ВКонтакте и Telegram)",
                  is_italic=True)
    add_table(document, title="Топ-10 просматриваемых публикаций", last_colum="Количество просмотров",
              items=[{"title": "", "url": "", "post_count": ""}] * 10)
    new_paragraph(document,
                  "ТОП-10 публикаций по сумме комментариев в социальных сетях (комментарии можно оставлять в ВКонтакте, Facebook, Twitter и Telegram)",
                  is_italic=True)
    add_table(document, title="Топ-10 обсуждаемых публикаций в социальных сетях", last_colum="Количество комментариев",
              items=[{"title": "", "url": "", "post_count": ""}] * 10)

    new_paragraph(document, "Половозрастная категория активных пользователей по теме", is_italic=True)

    group_1 = 0
    group_2 = 0
    group_3 = 0
    group_4 = 0
    for g in ages['graph_data']['group1']['graph']:
        group_1 += g[-1]
    for g in ages['graph_data']['group2']['graph']:
        group_2 += g[-1]
    for g in ages['graph_data']['group3']['graph']:
        group_3 += g[-1]
    for g in ages['graph_data']['group4']['graph']:
        group_4 += g[-1]
    new_paragraph(document,
                  f"""Наиболее активной аудиторией являются пользователи в возрасте"""
                  )
    add_pie_age(document, group_1, group_2, group_3, group_4)
    sex_m = 0
    sex_w = 0
    for m in ages['additional_data']['sex']['m']:
        sex_m += m[-1]
    for w in ages['additional_data']['sex']['w']:
        sex_w += w[-1]
    res = "женщины"
    if sex_w < sex_m:
        res = "мужчины"

    new_paragraph(document,
                  f"""Наиболее активной аудиторией с небольшим перевесом являются {res}."""
                  )

    add_pie_sex(document, sex_m, sex_w)

    spb = 0
    another = 0
    for c in city:
        if c['city'] == 'Санкт-Петербург':
            spb += int(c['users'])
        else:
            another += int(c['users'])
    res_active_user  = 0 if another + spb == 0 else int((spb * 100) / (another + spb))
    new_paragraph(document,
                  f"""Около {res}% активной аудитории располагается в Санкт-Петербурге."""
                  )
    add_pie_city(document, spb, another)

    new_paragraph(document, "Вывод", is_italic=True)

    new_paragraph(document,
                  f"""Всего за анализируемый период было размещено {get_str_int(res_dict["gs"]["total"])} публикаций в СМИ и {get_str_int(count)} публикаций в социальных сетях с упоминаниями по теме """)

    d_td_ = datetime.datetime.strptime(max_date_p, "%Y-%m-%d")
    new_paragraph(document, "Пиковый день по количеству публикаций приходятся"
                            f""" на {d_td_.day} {DATES[d_td_.month]} {d_td_.year} года. Всего в эту дату было опубликовано {get_str_int(smi[max_date_p] + social[max_date_p])} материала (из них {get_str_int(smi[max_date_p])} в СМИ и {get_str_int(social[max_date_p])}  """
                            "в социальных сетях). Резонансные темы этого дня следующие: ")

    new_paragraph(document, "По количеству просмотров лидирует . Темы публикаций с наибольшим количеством просмотров: ")
    date_t = datetime.datetime.strptime(cat[max_date], "%Y-%m-%d")
    new_paragraph(document,
                  f"""Пиковым днем по количеству лайков, репостов и комментариев является {date_t.day} {DATES[date_t.month]} {date_t.year} г. – {get_str_int(like_[max_date])} лайка, {get_str_int(repost_[max_date])} репоста, {get_str_int(comment_[max_date])} комментариев. Темы публикаций с наибольшим количеством реакций пользователей: """)
    smi_text = ""
    k = 0
    for i in top_smi["items"]:
        if k == 1:
            smi_text += ", "
        smi_text += "«" + i['title'] + "»"
        k += 1
        if k > 2:
            break
    social_text = ""
    k = 0
    for i in top_social["items"]:
        if k == 1:
            social_text += ", "
        social_text += "«" + i['title'] + "»"
        k += 1
        if k > 2:
            break
    new_paragraph(document,
                  f"""Наиболее активными источниками СМИ (информационные сайты) являются {smi_text},. Активными источниками социальных сетей стали сообщества {social_text}.""")

    new_paragraph(document,
                  f"""Наиболее активной аудиторией, отреагировавшей на публикации отчетного периода на тему, по возрастному критерию являются пользователи. По гендерному критерию более активными являются {res} с небольшим перевесом. {res_active_user}% активной аудитории располагается в Санкт-Петербурге.
 """)
    p = new_paragraph(document,
                      f"""Приложение 1. ТОП-10 публикаций за сутки за период {from_date.day}  по {to_date.day} {DATES[to_date.month]} {to_date.year}   года""")
    p.runs[-1].bold = True
    add_table_pril(document, top_post)
    return document


def new_paragraph(document, text=None, is_italic=False, paragraph=None, is_JUSTIFY=False):
    if paragraph is None:
        p = document.add_paragraph()
    else:
        p = document.paragraphs[paragraph]
    p.paragraph_format.first_line_indent = Mm(10)
    if is_JUSTIFY:
        p.alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.JUSTIFY

    if text:
        p.add_run(text, style=STYLE)
        if is_italic:
            p.runs[-1].italic = True
    return p


def add_title_text(document, text, is_bold=False, is_italic=False):
    parag_title = document.paragraphs[0]
    parag_title.add_run(
        text,
        style=STYLE
    )
    parag_title.alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER
    if is_bold:
        parag_title.runs[-1].bold = True
    if is_italic:
        parag_title.runs[-1].italic = True


if __name__ == '__main__':
    thread_id = 998
    _from = "2022-10-10 00:00:00"
    _to = "2022-10-17 00:00:00"
    login_user = "java_api"
    password = "4yEcwVnjEH7D"
    asyncio.run(prepare_report(thread_id, _from, _to, login_user, password))
