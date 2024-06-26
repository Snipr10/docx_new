import asyncio
import dateutil
import docx
import httpx as httpx
import traceback

import uvicorn
from dateutil.parser import parse
from datetime import datetime
from docx.oxml.shared import OxmlElement, qn
from docx.oxml.ns import nsdecls
from docx.oxml import parse_xml
from docx.enum.style import WD_STYLE_TYPE
from docx.shared import Inches, Pt
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.enum.table import WD_TABLE_ALIGNMENT
from datetime import timedelta
from dateutil.relativedelta import relativedelta
from io import BytesIO
import logging

from docx import Document
from fastapi import FastAPI, Request
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt, Inches
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from starlette.responses import StreamingResponse

from resp import post
from tonal_media import docx_tonal
from word_media import docx_media, login, convert_date, get_posts_info
from logging.config import dictConfig
from log_conf import log_config
from settings import SUBECT_URL, SUBECT_TOPIC_URL, STATISTIC_URL, STATISTIC_TRUST_GRAPH, GET_TRUST_URL, NETWORK_IDS

COOKIES = []

KOM_NAME = "Комитет по образованию"
STYLE = "Times New Roman"
PT = Pt(10.5)

dictConfig(log_config)
app = FastAPI()

UTC = 3

TIMEOUT = 7 * 60
# logger = logging.getLogger('foo-logger')

import logging

import rollbar
from rollbar.contrib.fastapi import LoggerMiddleware
from rollbar.logger import RollbarHandler

# Initialize Rollbar SDK with your server-side access token
rollbar.init(
    'd11e000ae6694189bfb39896cc4bcb6f',
    environment='staging',
    handler='async',
)

# Set root logger to log DEBUG and above
logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG)
rollbar_handler = RollbarHandler()
rollbar_handler.setLevel(logging.ERROR)
logger.addHandler(rollbar_handler)
app.add_middleware(LoggerMiddleware)  # should be added as the last middleware


@app.post('/get_report')
async def index(request: Request):
    body_json = await request.json()

    period = body_json.get('period', None)
    _from_data = body_json.get('from', None)
    _to_data = body_json.get('to', None)

    periods_data = {"period": period, "_from_data": _from_data, "_to_data": _to_data}
    reference_ids_str = body_json.get('reference_ids')

    reference_ids = []

    for id_ in reference_ids_str:
        reference_ids.append(int(id_))

    logger.info(f"body_json {body_json}")

    attempt = 0
    max = 3
    while attempt < max:
        try:
            document = await creater(reference_ids, body_json.get('login'), body_json.get('password'),
                                     int(body_json.get('thread_id')), periods_data)
            f = BytesIO()

            document.save(f)
            f.seek(0)

            response = StreamingResponse(f, media_type="text/docx")
            response.headers["Content-Disposition"] = "attachment; filename=report.docx"
            attempt = max
            return response
        except Exception as e:
            logger.error(f'get_report {e}')
            attempt += 1
    return "Что-то пошло не так"


@app.post('/get_new_report')
async def new_report(request: Request):
    from new_report.new_report import prepare_report

    body_json = await request.json()

    _from = body_json.get('from')
    _to = body_json.get('to')
    thread_id = int(body_json.get('thread_id'))
    attempt = 0
    max = 3
    while attempt < max:
        try:
            document = await prepare_report(thread_id, _from, _to, body_json.get('login'), body_json.get('password'))
            f = BytesIO()

            document.save(f)
            f.seek(0)

            response = StreamingResponse(f, media_type="text/docx")
            response.headers["Content-Disposition"] = "attachment; filename=report.docx"
            attempt = max
            return response
        except Exception as e:
            logger.error(f'get_new_report {e}')
            attempt += 1
    return "Что-то пошло не так"


@app.post('/get_publication_summary')
async def index_media(request: Request):
    body_json = await request.json()

    _from = body_json.get('from')
    _to = body_json.get('to')

    _sort = bool(body_json.get('sort', 0))

    referenceFilter = []
    for id_ in body_json.get('reference_ids'):
        referenceFilter.append(int(id_))

    network_id = []
    for id_ in body_json.get('network_id'):
        network_id.append(int(id_))
    thread_id = int(body_json.get('thread_id'))
    if not network_id:
        network_id = NETWORK_IDS
    friendly_ids = body_json.get('friendly', [])
    trustoption = body_json.get('trustoption', None)
    try:
        document = await docx_media(thread_id, _from, _to,
                                    referenceFilter, network_id, body_json.get('user_id'), friendly_ids, trustoption, _sort)

        f = BytesIO()
        document.save(f)
        f.seek(0)

        response = StreamingResponse(f, media_type="text/docx")
        response.headers["Content-Disposition"] = "attachment; filename=report.docx"
        return response
    except Exception as e:
        logger.error(f'get_publication_summary {e}')
        return "Что-то пошло не так"


@app.post('/get_report/tonal')
async def tonal(request: Request):
    body_json = await request.json()

    _from = body_json.get('from')
    _to = body_json.get('to')

    _from_parse = convert_date(_from)
    _to_parse = convert_date(_to)

    thread_ids_str = body_json.get('thread_ids')
    iogv_name = body_json.get('iogv_name')
    types = body_json.get('type')
    smi_type = body_json.get('smi_type', 'any')

    thread_ids = []

    for id_ in thread_ids_str:
        thread_ids.append(int(id_))

    try:
        document = await docx_tonal(thread_ids, _from_parse, _to_parse, iogv_name, types, smi_type)

        f = BytesIO()
        document.save(f)
        f.seek(0)

        response = StreamingResponse(f, media_type="text/docx")
        response.headers["Content-Disposition"] = "attachment; filename=report.docx"
        return response
    except Exception as e:
        logger.error(f'get_report/tonal {e}')
        return "Что-то пошло не так"


def _last_time(day):
    return datetime(day.year, day.month, day.day, 23, 59, 59).strftime('%Y-%m-%d %H:%M:%S')


async def creater(reference_ids, login_user, password, thread_id, periods_data):
    async with httpx.AsyncClient() as session:

        await login(session, login_user, password)

        today_all = datetime.today() + timedelta(hours=UTC)
        logger.error(f"today_all")

        if periods_data.get("period") == "day":
            today_all = datetime.today() + timedelta(hours=UTC)
            today = today_all.strftime('%d-%m-%Y')
            today_str = f"на {today}"
            periods_data["_from_data"] = get_from_date(periods_data.get("period"))
            periods_data["_to_data"] = _last_time(today_all)

        else:
            if periods_data.get("period"):
                today_all = datetime.today() + timedelta(hours=UTC) - timedelta(days=1)
                today_all = datetime(today_all.year, today_all.month, today_all.day, 23, 59, 59)
                today = today_all.strftime('%d-%m-%Y')
                today_str = f"за период с {get_from_date_datetime(periods_data.get('period')).strftime('%d-%m-%Y')} по {today}"
                periods_data["_from_data"] = get_from_date(periods_data.get("period"))
                periods_data["_to_data"] = today_all.strftime('%Y-%m-%d %H:%M:%S')
            else:
                _to_data = dateutil.parser.parse(periods_data["_to_data"])
                today_all = datetime(_to_data.year, _to_data.month, _to_data.day, 23, 59, 59)
                today_str = f"за период с {dateutil.parser.parse(periods_data['_from_data']).strftime('%d-%m-%Y')} по {dateutil.parser.parse(periods_data['_to_data']).strftime('%d-%m-%Y')}"

                periods_data["_to_data"] = today_all.strftime('%Y-%m-%d %H:%M:%S')
        logger.error(f"document")

        document = Document()

        obj_styles = document.styles
        obj_charstyle = obj_styles.add_style(STYLE, WD_STYLE_TYPE.CHARACTER)
        obj_font = obj_charstyle.font
        obj_font.size = Pt(10.5)
        obj_font.name = STYLE
        logger.error(f"add_title")

        add_title(document, today_str)
        logger.error(f"sub")

        sub = await get_start_date(session)
        logger.error(f"try sub")

        try:

            topics_tables, statistic_tables, trust_tables, charts_data, posts_info = await get_tables(session, periods_data, sub,
                                                                                          thread_id,
                                                                                          reference_ids)
        except Exception as e:
            logger.error(f'creater {e}')

            try:
                topics_tables.cancel()
            except Exception as e:
                logger.error(f"topics_tables {e}")
            try:
                statistic_tables.cancel()
            except Exception as e:
                logger.error(f"statistic_tables {e}")
            try:
                trust_tables.cancel()
            except Exception as e:
                logger.error(f"trust_tables {e}")
            try:
                charts_data.cancel()
            except Exception as e:
                logger.error(f"charts_data {e}")
            raise e

        table_number = 1

        add_table_title = True
        for topics_table_title, topics_table_data, reference_id in topics_tables:
            if add_table_title:
                add_title_text(document, "Главные темы публикаций в СМИ", True)
            add_table1(document, table_number, topics_table_title, topics_table_data, today_str, add_table_title, posts_info)
            table_number += 1
            add_table_title = False

        add_table_title = True
        for statistic_table_title, statistic_table_date in statistic_tables:
            if add_table_title:
                document.add_page_break()
                add_title_text(document, "\n Статистика по публикациям с упоминанием субъектов", True)
            add_table2(document, table_number, statistic_table_date, statistic_table_title, today_str, add_table_title, posts_info)
            table_number += 1
            add_table_title = False

        chart_number = 1
        add_chart_title = True
        for statistic_chart_title, statist_chart_data in charts_data:
            if add_chart_title:
                document.add_page_break()
                add_title_text(document, "Динамика распространения публикаций", True)
                add_chart_title = False
            chart_number = add_chart_document(document, chart_number, statistic_chart_title, statist_chart_data,
                                              today_str,
                                              today_all,
                                              periods_data)
            chart_number += 1

        if chart_number % 2 == 0:
            document.add_page_break()

        add_title_text(document, "ТОПы публикаций СМИ и социальных сетей", True)

        first = True
        for trust_table_title, table_social_data_range, table_smi_data_range, table_social_data_pos_neu, \
            table_smi_data_pos_neu, table_social_data_neg, table_smi_data_neg in trust_tables:

            if table_smi_data_range or table_smi_data_pos_neu or table_smi_data_neg:
                try:
                    add_table_trust(
                        document,
                        table_number,
                        trust_table_title,
                        table_smi_data_range,
                        table_smi_data_pos_neu,
                        table_smi_data_neg,
                        today_str,
                        "СМИ",
                        first
                    )
                    first = False
                    table_number += 1
                    document.add_page_break()
                except Exception as e:
                    logger.error(f"creater 1 {e}")

            if table_social_data_range or table_social_data_pos_neu or table_social_data_neg:
                try:

                    add_table_trust(
                        document,
                        table_number,
                        trust_table_title,
                        table_social_data_range,
                        table_social_data_pos_neu,
                        table_social_data_neg,
                        today_str,
                        "в социальных сетях",
                        first, True
                    )
                    first = False
                    table_number += 1
                    document.add_page_break()
                except Exception as e:
                    logger.error(f"creater 2 {e}")

        return document


def change_table_font(table):
    for row in table.rows:
        for cell in row.cells:
            paragraphs = cell.paragraphs
            update_pagagraphs(paragraphs)


def update_pagagraphs(paragraphs):
    for paragraph in paragraphs:
        for run in paragraph.runs:
            font = run.font
            font.size = Pt(10.5)
            font.name = STYLE


def add_title(document, today):
    add_title_text(document,
                   f'Отчет по публикациям в Личном Кабинете Мониторинговой системы {today},'
                   f' созданный на основании публикаций СМИ и социальных сетей',
                   False
                   )


def add_title_text(document, text, is_bold):
    parag_title = document.add_paragraph()
    parag_title.add_run(
        text,
        style=STYLE
    )
    parag_title.alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER
    if is_bold:
        parag_title.runs[-1].bold = True


def set_cell_vertical_alignment(cell, align="center"):
    try:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        tcValign = OxmlElement('w:vAlign')
        tcValign.set(qn('w:val'), align)
        tcPr.append(tcValign)
        return True
    except Exception:
        traceback.print_exc()
        return False


def set_center(cell):
    set_cell_vertical_alignment(cell)
    cell.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell.paragraphs[0].paragraph_format.alignment = WD_TABLE_ALIGNMENT.CENTER


def set_right(cell):
    set_cell_vertical_alignment(cell)
    cell.alignment = WD_TABLE_ALIGNMENT.RIGHT
    cell.paragraphs[0].paragraph_format.alignment = WD_TABLE_ALIGNMENT.RIGHT


def set_left(cell):
    set_cell_vertical_alignment(cell)
    cell.alignment = WD_TABLE_ALIGNMENT.LEFT
    cell.paragraphs[0].paragraph_format.alignment = WD_TABLE_ALIGNMENT.LEFT


def add_table1(document, table_number, header, records, today, add_table_title, posts_info):
    parag_table_1 = document.add_paragraph()
    text = f' Таблица {table_number} - Главные темы публикаций СМИ с упоминаниями '
    if not add_table_title:
        text = "\n" + text
    parag_table_1.add_run(
        text,
        style=STYLE
    )
    add_name(parag_table_1, header)
    parag_table_1.add_run(
        f' {today}',
        style=STYLE
    )
    parag_table_1.paragraph_format.space_after = Inches(0)
    parag_table_1.alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.JUSTIFY

    table = document.add_table(rows=0, cols=6)
    table.autofit = False
    table.allow_autofit = False
    table.columns[0].width = Inches(0.3)
    table.columns[1].width = Inches(2.45)
    table.columns[2].width = Inches(1.0)
    table.columns[3].width = Inches(0.7)
    table.columns[4].width = Inches(1.0)
    table.columns[5].width = Inches(0.7)

    table.style = 'TableGrid'

    hdr_cells = table.add_row().cells

    hdr_cells[0].text = " "
    hdr_cells[1].text = "Тема"
    hdr_cells[2].text = "Публикаций с упоминанием субъекта"
    hdr_cells[3].text = "Охват публикаций с упоминанием"
    hdr_cells[4].text = "Всего публикаций в теме"
    hdr_cells[5].text = "Охват всех публикаций"

    set_center(hdr_cells[2])
    set_center(hdr_cells[3])
    set_center(hdr_cells[4])
    set_center(hdr_cells[5])
    set_cell_vertical_alignment(hdr_cells[1])

    i = 1
    max_count = 0
    for cell in records:
        if max_count >= 20:
            break
        max_count += 1
        row_cells = table.add_row().cells
        row_cells[5].text = str(cell['total_attendance'])
        row_cells[4].text = str(cell['total_posts'])
        row_cells[3].text = str(cell['attendance'])
        row_cells[2].text = str(cell['postcount'])

        row_cells[1].text = str(cell['title'])
        row_cells[0].text = str(i)
        row_cells[0].alignment = WD_TABLE_ALIGNMENT.RIGHT
        row_cells[0].paragraphs[0].paragraph_format.alignment = WD_TABLE_ALIGNMENT.RIGHT

        set_center(row_cells[2])
        set_center(row_cells[3])
        set_center(row_cells[4])
        set_center(row_cells[5])

        i += 1
    change_table_font(table)


def add_table2(document, table_number, records, table_type, today, add_table_title, posts_info):
    parag_table = document.add_paragraph()
    text = f' Таблица {table_number}  - Общая статистика публикаций {table_type} с упоминаниями субъектов {today}'
    if not add_table_title:
        text = "\n" + text
    parag_table.add_run(
        text,
        style=STYLE
    )

    parag_table.paragraph_format.space_after = Inches(0)
    parag_table.alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.JUSTIFY

    table = document.add_table(rows=1, cols=6)
    table.style = 'TableGrid'
    table.columns[0].width = Inches(2.25)
    table.columns[1].width = Inches(0.8)
    table.columns[2].width = Inches(0.8)
    table.columns[3].width = Inches(0.8)
    table.columns[4].width = Inches(0.8)
    table.columns[5].width = Inches(0.8)

    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = "Субъекты"
    set_cell_vertical_alignment(hdr_cells[0])

    hdr_cells[1].text = "Количество публикаций, всего"
    set_center(hdr_cells[1])
    table.rows[0].width = Inches(2)
    hdr_cells[2].text = "Дружественные"
    set_center(hdr_cells[2])
    hdr_cells[3].text = "Позитивные публикации"
    set_center(hdr_cells[3])
    hdr_cells[4].text = "Негативные публикации"
    set_center(hdr_cells[4])
    hdr_cells[5].text = "Нейтральные публикации"
    set_center(hdr_cells[5])
    for cell in records:
        row_cells = table.add_row().cells

        row_cells[0].text = cell['header']
        positive = int(cell['positive']['posts'])
        negative = int(cell['negative']['posts'])
        netural = int(cell['netural']['posts'])
        total = int(cell['total']['posts'])

        row_cells[1].text = str(total)
        set_center(row_cells[1])
        _, _, _, _, friendly_smi, friendly_social = posts_info[cell["reference_id"]]
        friendly = 0
        if table_type == "СМИ":
            friendly = friendly_smi
        else:
            friendly = friendly_social
        row_cells[2].text = str(friendly)
        set_center(row_cells[2])
        row_cells[3].text = str(positive)
        set_center(row_cells[3])
        row_cells[4].text = str(negative)
        set_center(row_cells[4])

        row_cells[5].text = str(total - positive - negative)
        set_center(row_cells[5])

    change_table_font(table)


def update_center_right(row_cell):
    set_cell_vertical_alignment(row_cell)
    row_cell.paragraphs[0].alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.RIGHT


async def subects(session):
    response = await post(session, SUBECT_URL, None)
    try:
        res = []
        for r in response.json():
            res.extend(r['items'] or [])
        return res
    except Exception as e:
        logger.error(f"subects {e}")
        return []


async def subects_topic(session, reference_id, thread_id, periods_data, table_name):
    try:
        if periods_data.get("period"):
            payload = {
                "thread_id": thread_id,
                "referenceFilter": [reference_id],
                "period": periods_data.get("period"),
                "type": "smi",
                "start": 0,
                "limit": 100
            }
        else:
            payload = {
                "thread_id": thread_id,
                "referenceFilter": [reference_id],
                "from": periods_data.get("_from_data"),
                "to": periods_data.get("_to_data"),
                "type": "smi",
                "start": 0,
                "limit": 100
            }
        response = await post(session, SUBECT_TOPIC_URL, payload)

        res = []
        try:
            for r in response.json().get("items", []):
                res.append(r)
        except Exception as e:
            logger.error(f"subects_topic {e} {response.text}")
        return res, table_name, reference_id
    except Exception as e:
        logger.error(f"subects_topic {e}")
        raise e


def get_from_date(period):
    date_from = get_from_date_datetime(period)

    return datetime(date_from.year, date_from.month, date_from.day).strftime('%Y-%m-%d %H:%M:%S')


def get_from_date_datetime(period):
    date_from = None
    if period == "day":
        date_from = datetime.today() + timedelta(hours=UTC)
    elif period == "week":
        date_from = datetime.today() - timedelta(days=7) + timedelta(hours=UTC)
    elif period == "month":
        date_from = datetime.today() - relativedelta(months=1) + timedelta(hours=UTC)
    return date_from


async def subects_static(session, reference_id, thread_id, periods_data, table_name):
    try:
        payload = {
            "thread_id": thread_id,
            "from": periods_data.get("_from_data"),
            "to": periods_data.get("_to_data"),
            "filter": {"referenceFilter": [reference_id]}
        }
        response = await post(session, STATISTIC_URL, payload)

        res_gs = {}
        res_soc = {}
        keys = ["fb", "vk", "tw", "tg", "ig", "yt"]
        try:
            res = response.json()
            if res.get("gs", {}).get("total", {}).get("posts", 0) is not None and res.get("gs", {}).get("total",
                                                                                                        {}).get(
                "posts", 0) > 0:
                res_gs = response.json().get("gs", {})
            total_posts = 0
            total_positive = 0
            total_negative = 0
            total_netural = 0
            for k in keys:
                total_posts += res[k]['total']['posts']
                total_positive += res[k]['positive']['posts']
                total_negative += res[k]['negative']['posts']
                total_netural += res[k]['netural']['posts']

            social = {
                'total': {
                    'posts': total_posts
                },
                'positive': {
                    'posts': total_positive
                },
                'negative': {
                    'posts': total_negative
                },
                'netural': {
                    'posts': total_netural
                }
            }
            if social.get("total", {}).get("posts", 0) > 0:
                res_soc = social
        except Exception as e:
            logger.error(f"subects_static {e} {response.text}")
        return res_gs, res_soc, table_name, reference_id
    except Exception as e:
        logger.error(f"subects_static {e}")
        raise e


async def add_topics(session, periods_data, sub, thread_id, reference_ids):
    async with httpx.AsyncClient(cookies=session.cookies) as session:
        try:
            tables = []
            table_gather = []
            for s in sub:
                table_name = s['keyword']
                reference_id = s['id']
                if reference_id in reference_ids:
                    table_gather.append(subects_topic(session, reference_id, thread_id, periods_data, table_name))
            for table_data, table_name, reference_id in await asyncio.gather(*table_gather):
                if table_data:
                    tables.append((table_name, table_data, reference_id))
            return tables
        except Exception as e:
            logger.error(f"add_topics {e}")
            raise e


async def add_statistic(session, periods_data, sub, thread_id, reference_ids):
    async with httpx.AsyncClient(cookies=session.cookies) as session:
        try:
            tables = []
            table_data_smi = []
            table_data_soc = []
            table_gather = []
            for s in sub:
                reference_id = s['id']
                if reference_id in reference_ids:
                    table_gather.append(subects_static(session, reference_id, thread_id, periods_data, s['keyword']))
            for row_gs, ros_soc, table_name, reference_id in await asyncio.gather(*table_gather):
                if row_gs:
                    row_gs["header"] = table_name
                    row_gs["reference_id"] = reference_id

                    table_data_smi.append(row_gs)
                if ros_soc:
                    ros_soc["header"] = table_name
                    ros_soc["reference_id"] = reference_id
                    table_data_soc.append(ros_soc)
            if table_data_smi:
                tables.append(("СМИ", table_data_smi))
            if table_data_soc:
                tables.append(("в социальных сетях", table_data_soc))
            return tables
        except Exception as e:
            logger.error(f"add_statistic {e}")
            raise e


async def get_trust_stat(session, thread_id, reference_ids, periods_data, network_id, post_count, negative=None):
    try:
        payload = {
            "thread_id": thread_id,
            "negative": negative,
            "post_count": post_count,
            "from": periods_data.get("_from_data"),
            "to": periods_data.get("_to_data"),
            "filter": {"network_id": network_id, "referenceFilter": [reference_ids]}
        }
        response = await post(session, GET_TRUST_URL, payload)

        if response.status_code == 405:
            return []
        return response.json()
    except Exception as e:
        logger.error(f"get_trust_stat {e}")
        raise e


def add_table_trust(document, table_number, header, table_data_range,
                    table_data_pos_neu,
                    table_data_neg, today, doc_type, first, social=False):
    parag_table_1 = document.add_paragraph()
    p_text = f' Таблица {table_number} - ТОПы публикаций {doc_type} с упоминаниями '

    if not first:
        p_text = "\n " + p_text

    parag_table_1.add_run(
        p_text,
        style=STYLE
    )

    add_name(parag_table_1, header)
    parag_table_1.add_run(
        f' {today}.',
        style=STYLE
    )

    p_small_text = "\nФормирование ТОП публикаций осуществляется на основании охватов издания." if \
        not social else "\nФормирование ТОП публикаций осуществляется на основании суммы реакций."

    parag_table_1.add_run(
        p_small_text,
        style=STYLE
    )
    parag_table_1.runs[-1].font.size = Pt(8)

    parag_table_1.paragraph_format.space_after = Inches(0)
    parag_table_1.alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.JUSTIFY

    table = document.add_table(rows=0, cols=4)
    table.autofit = False
    table.allow_autofit = False
    if not social:
        table.columns[0].width = Inches(0.2)
        table.columns[1].width = Inches(1.5)
        table.columns[2].width = Inches(0.6)
        table.columns[3].width = Inches(3.85)
    else:
        table.columns[0].width = Inches(0.2)
        table.columns[1].width = Inches(1.2)
        table.columns[2].width = Inches(1.5)
        table.columns[3].width = Inches(3.25)
    table.style = 'TableGrid'
    add_col_name(table, social)

    change_table_font(table)

    if table_data_range:
        table_data_range = sorted(table_data_range, key=lambda x: x[0], reverse=True)
        row_cells = table.add_row().cells
        if not social:
            header_cell(row_cells, "ТОП-5 публикаций по охватам", "81e5f8")
        else:
            header_cell(row_cells, "ТОП-5 публикаций", "81e5f8")
        add_top5(table, table_data_range, social)
    if table_data_pos_neu:
        table_data_pos_neu = sorted(table_data_pos_neu, key=lambda x: x[0], reverse=True)
        row_cells = table.add_row().cells
        header_cell(row_cells, "ТОП-5 позитивных и нейтральных публикаций", "72f983")
        add_top5(table, table_data_pos_neu, social)
    if table_data_neg:
        table_data_neg = sorted(table_data_neg, key=lambda x: x[0], reverse=True)
        row_cells = table.add_row().cells
        header_cell(row_cells, "ТОП-5 негативных и противоречивых публикаций", "d24141")
        add_top5(table, table_data_neg, social)


def add_col_name(table, social):
    row_cells = table.add_row().cells
    row_cells[1].text = "Ссылка"
    if social:
        row_cells[2].text = "Реакции"
    else:
        row_cells[2].text = "Охват"
    row_cells[3].text = "Текст"


def header_cell(hdr_cells, header, color):
    hdr_cells[0].text = header
    hdr_cells[0].merge(hdr_cells[1])
    hdr_cells[0].merge(hdr_cells[2])
    hdr_cells[0].merge(hdr_cells[3])
    shading_elm = parse_xml(r'<w:shd {} w:fill="{}"/>'.format(nsdecls('w'), color))
    hdr_cells[0]._tc.get_or_add_tcPr().append(shading_elm)


def add_top5(table, table_data, social):
    for i in range(len(table_data)):
        row_cells = table.add_row().cells
        try:
            row_cells[0].paragraphs[0].part.style = STYLE
            row_cells[2].paragraphs[0].part.style = STYLE
            row_cells[3].paragraphs[0].part.style = STYLE
            row_cells[0].text = str(i + 1)
            if social:
                row_cells[2].paragraphs[0].add_run(
                    "Просмотры: ",
                    style=STYLE
                )
                row_cells[2].paragraphs[0].runs[-1].bold = True

                row_cells[2].paragraphs[0].add_run(
                    table_data[i][1]["viewed"] + "\n",
                    style=STYLE
                )

                row_cells[2].paragraphs[0].add_run(
                    "Лайки: ",
                    style=STYLE
                )
                row_cells[2].paragraphs[0].runs[-1].bold = True

                row_cells[2].paragraphs[0].add_run(
                    table_data[i][1]["likes"] + "\n",
                    style=STYLE
                )
                row_cells[2].paragraphs[0].add_run(
                    "Комментарии: ",
                    style=STYLE
                )
                row_cells[2].paragraphs[0].runs[-1].bold = True

                row_cells[2].paragraphs[0].add_run(
                    table_data[i][1]["comments"] + "\n",
                    style=STYLE
                )
                row_cells[2].paragraphs[0].add_run(
                    "Репосты: ",
                    style=STYLE
                )
                row_cells[2].paragraphs[0].runs[-1].bold = True

                row_cells[2].paragraphs[0].add_run(
                    table_data[i][1]["reposts"] + "\n",
                    style=STYLE
                )
            else:
                row_cells[2].text = str(table_data[i][0]) + "\n"

            row_cells[3].paragraphs[0].add_run(
                f"{table_data[i][1]['created_date']}\n",
                style=STYLE
            )
            row_cells[3].paragraphs[0].runs[-1].italic = True
            row_cells[3].paragraphs[0].runs[-1].bold = True
            row_cells[3].paragraphs[0].runs[-1].font.size = Pt(8)

            if table_data[0][1]['title']:
                row_cells[3].paragraphs[0].add_run(
                    f"{table_data[i][1]['title']}\n\n",
                    style=STYLE
                )
                row_cells[3].paragraphs[0].runs[-1].bold = True
            else:
                row_cells[3].paragraphs[0].add_run(
                    "\n",
                    style=STYLE
                )
            add_hyperlink(row_cells[1].paragraphs[0], table_data[i][1]['url'] or " ", table_data[i][1]['url'] or " ", None, True)

            text, add_link = remove_html_tags(table_data[i][1]['text'] or " ")
            row_cells[3].paragraphs[0].add_run(text)
            if add_link:
                add_hyperlink(row_cells[3].paragraphs[0], table_data[i][1]['url'] or " ", "далее по ссылке", None, True, True)

            set_center(row_cells[2])
        except Exception as e:
            logger.error(f"add_top5 {e}")


def add_to5_title(table, title):
    row_cells = table.add_row().cells
    row_cells[0].text = title
    row_cells[0].merge(row_cells[1])
    set_cell_vertical_alignment(row_cells[0])
    row_cells[0].paragraphs[0].alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.CENTER


def get_text(dict_s, main_text, is_all=True):
    text = dict_s[1]['author'] + "\n" + dict_s[1]['created_date'] + "\n" + remove_html_tags(dict_s[1][main_text])
    if not is_all:
        text = text[: 100]
    return text


def remove_html_tags(text, len=200):
    add_link = False
    try:
        import re
        clean = re.compile('<.*?>')
        update_text = re.sub(clean, '', text)
        if update_text.__len__() > 200:
            update_text = update_text[:len - 19] + "... "
            add_link = True
        return update_text, add_link
    except Exception as e:
        logger.error(f"remove_html_tags {e}")

        return '', add_link


async def get_attendance(session, res_net_social):
    try:
        res_gather = []
        for r in res_net_social:
            res_gather.append(get_attendance_data(session, r))
        res = await asyncio.gather(*res_gather)

        return sorted(res, key=lambda x: x[0], reverse=True)[:5]
    except Exception as e:
        logger.error(f"get_attendance {e}")
        raise e


async def get_attendance_data(session, r):
    return r['supercoefficient'], {
        "created_date": r['created_date'],
        "author": r['author'],
        "text": r["text"],
        "title": r["title"],
        "url": r["uri"],
        "likes": r.get("likes", 0),
        "reposts": r.get("reposts", 0),
        "comments": r.get("comments", 0),
        "viewed": r.get("viewed", 0)

    }


async def get_trust(session, periods_data, sub, thread_id, reference_ids):
    async with httpx.AsyncClient(cookies=session.cookies) as session:
        try:
            tables = []

            network_ids = NETWORK_IDS

            table_gather = []

            for s in sub:
                reference_id = s['id']
                if reference_id in reference_ids:
                    table_gather.append(
                        get_trust_for_sub(session, reference_id, network_ids, s['keyword'], periods_data, thread_id))
            for trust_state_date in await asyncio.gather(*table_gather):
                if trust_state_date is not None:
                    tables.append(trust_state_date)
            return tables
        except Exception as e:
            logger.error(f"get_trust {e}")
            raise e


async def get_trust_res_net_social_range(session, network_ids, thread_id, reference_id, periods_data):
    try:
        res_net_social_range_gather = []
        res_net_social_range = []

        for net_id in network_ids:
            res_net_social_range_gather.append(
                get_trust_stat(session, thread_id, reference_id, periods_data, [net_id], 3, None))

        for trust_state_date in await asyncio.gather(*res_net_social_range_gather):
            res_net_social_range.extend(trust_state_date)

        return res_net_social_range
    except Exception as e:
        logger.error(f"get_trust_res_net_social_range {e}")
        raise e


async def get_trust_for_sub(session, reference_id, network_ids, title, periods_data, thread_id):
    try:
        table = None
        networks_without_g = network_ids.copy()
        networks_without_g.remove(4)

        res_net_social_range, res_net_gs_range, res_net_social_pos_neu, res_net_gs_range_pos_neu, res_net_social_neg, res_net_gs_range_neg = await asyncio.gather(
            get_trust_res_net_social_range(session, networks_without_g, thread_id, reference_id, periods_data),
            get_trust_stat(session, thread_id, reference_id, periods_data, [4], 5, None),
            get_trust_stat(session, thread_id, reference_id, periods_data, networks_without_g, 5, False),
            get_trust_stat(session, thread_id, reference_id, periods_data, [4], 5, False),
            get_trust_stat(session, thread_id, reference_id, periods_data, networks_without_g, 5, True),
            get_trust_stat(session, thread_id, reference_id, periods_data, [4], 5, True)
        )

        if len(res_net_social_range) > 0 or len(res_net_gs_range) > 0 or len(res_net_social_pos_neu) > 0 or len(
                res_net_gs_range_pos_neu) > 0 or len(res_net_social_neg) > 0 or len(res_net_gs_range_neg) > 0:
            table_social_data_range, table_smi_data_range, table_social_data_pos_neu, table_smi_data_pos_neu, table_social_data_neg, table_smi_data_neg = await asyncio.gather(
                get_attendance(session, res_net_social_range),
                get_attendance(session, res_net_gs_range),
                get_attendance(session, res_net_social_pos_neu),
                get_attendance(session, res_net_gs_range_pos_neu),
                get_attendance(session, res_net_social_neg),
                get_attendance(session, res_net_gs_range_neg)
            )
            table = (title,
                     table_social_data_range, table_smi_data_range,
                     table_social_data_pos_neu, table_smi_data_pos_neu,
                     table_social_data_neg, table_smi_data_neg)
        return table
    except Exception as e:
        logger.error(f"get_trust_for_sub {e}")
        raise e


async def get_start_date(session):
    return await subects(session)


# async def get_posts_statistic(session, periods_data, sub, thread_id, reference_ids):
#     async with httpx.AsyncClient(cookies=session.cookies) as session:
#         try:
#             tables = []
#             table_gather = []
#             for s in sub:
#                 chart_name = s['keyword']
#                 reference_id = s['id']
#                 if reference_id in reference_ids:
#                     table_gather.append(post_static(session, reference_id, thread_id, periods_data, chart_name))
#             for table_data, chart_name in await asyncio.gather(*table_gather):
#                 if table_data:
#                     tables.append((chart_name, table_data))
#             return tables
#         except Exception as e:
#             logger.error(f"get_posts_statistic {e}")
#             raise e

async def get_posts_statistic(session, periods_data, sub, thread_id, reference_ids):
    async with httpx.AsyncClient(cookies=session.cookies) as session:
        try:
            tables = []
            table_gather = []
            for s in sub:
                chart_name = s['keyword']
                reference_id = s['id']
                if reference_id in reference_ids:
                    table_gather.append(post_static(session, reference_id, thread_id, periods_data, chart_name))
            for table_data, chart_name in await asyncio.gather(*table_gather):
                if table_data:
                    tables.append((chart_name, table_data))
            return tables
        except Exception as e:
            logger.error(f"get_posts_statistic {e}")
            raise e


# async def post_static(session, reference_id, thread_id, periods_data, chart_name):
#     limit = 200
#     start = 0
#     posts = []
#     while True:
#         payload = {
#             "thread_id": thread_id,
#             "from": periods_data.get("_from_data"),
#             "to": periods_data.get("_to_data"),
#             "limit": limit, "start": start, "sort": {"type": "date", "order": "desc", "name": "dateDown"},
#             "filter": {"network_id": [1, 2, 3, 4, 5, 7, 8],
#                        "referenceFilter": [reference_id], "repostoption": "whatever"}
#         }
#         response = await post(session, STATISTIC_POST_URL, payload)
#
#         posts.extend(response.json().get("posts") or [])
#         if not response.json().get("posts") or response.json().get("count") <= len(posts):
#             break
#         start += limit
#     return posts, chart_name


async def post_static(session, reference_id, thread_id, periods_data, chart_name):
    payload = {
        "thread_id": thread_id,
        "from": periods_data.get("_from_data"),
        "to": periods_data.get("_to_data"),
        "filter": {"network_id": NETWORK_IDS,
                   "referenceFilter": [reference_id]}
    }
    try:
        response = await post(session, STATISTIC_TRUST_GRAPH, payload)
    except Exception as e:
        pass
    return response.json(), chart_name


async def get_tables(session, periods_data, sub, thread_id, reference_ids):
    trust_tables, topics_tables, statistic_tables, charts_data, posts_info = await asyncio.gather(
        get_trust(session, periods_data, sub, thread_id, reference_ids),
        add_topics(session, periods_data, sub, thread_id, reference_ids),
        add_statistic(session, periods_data, sub, thread_id, reference_ids),
        get_posts_statistic(session, periods_data, sub, thread_id, reference_ids),
        get_posts_info(session, thread_id, periods_data, reference_ids)
    )

    return topics_tables, statistic_tables, trust_tables, charts_data, posts_info


def add_hyperlink(paragraph, url, text, color, underline, is_italic=False):
    part = paragraph.part
    r_id = part.relate_to(url, docx.opc.constants.RELATIONSHIP_TYPE.HYPERLINK, is_external=True)

    hyperlink = docx.oxml.shared.OxmlElement('w:hyperlink')
    hyperlink.set(docx.oxml.shared.qn('r:id'), r_id, )

    new_run = docx.oxml.shared.OxmlElement('w:r')

    rPr = docx.oxml.shared.OxmlElement('w:rPr')
    if is_italic:
        u = docx.oxml.shared.OxmlElement('w:i')
        u.append(docx.oxml.shared.OxmlElement('w:iCs'))
        rPr.append(u)
    if not color is None:
        c = docx.oxml.shared.OxmlElement('w:color')
        c.set(docx.oxml.shared.qn('w:val'), color)
        rPr.append(c)

    if not underline:
        u = docx.oxml.shared.OxmlElement('w:u')
        u.set(docx.oxml.shared.qn('w:val'), 'none')
        rPr.append(u)

    new_run.append(rPr)
    new_run.text = text
    new_run.style = STYLE
    hyperlink.append(new_run)
    paragraph._p.append(hyperlink)
    # paragraph._p.style = STYLE
    paragraph.style.font.name = "Times New Roman"
    paragraph.style.font.size = docx.shared.Pt(10.5)

    return hyperlink


def change_color(series, color):
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = color


def change_color_line(line, color):
    line.color.rgb = color


def add_name(p, name):
    p.add_run(
        name,
        style=STYLE
    )
    p.runs[-1].bold = True


def update_chart_none(data_list):
    for i in range(len(data_list)):
        if data_list[i] == 0:
            data_list[i] = ""
    return data_list


def update_chart_style(chart):
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    chart.legend.font.name = STYLE

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(5)

    value_axis = chart.value_axis
    value_axis.has_major_gridlines = False

    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.font.name = STYLE
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.name = STYLE
    chart.plots[0].chart.value_axis.tick_labels.font.size = Pt(10)
    chart.plots[0].chart.value_axis.tick_labels.font.name = STYLE
    chart.plots[0].chart.category_axis.tick_labels.font.size = Pt(10)
    chart.plots[0].chart.category_axis.tick_labels.font.name = STYLE

    shape_properties = OxmlElement("c:spPr")
    chart.element.append(shape_properties)

    fill_properties = OxmlElement("a:ln")
    shape_properties.append(fill_properties)
    scheme_color = OxmlElement("a:noFill")

    fill_properties.append(scheme_color)


def add_chart_document(document, chart_number, statistic_chart_title, statist_chart_data, today, today_all,
                       periods_data):
    parag_table = document.add_paragraph()
    parag_table.add_run(
        f' График {chart_number} - Динамика распространения публикаций с упоминанием ',
        style=STYLE
    )
    add_name(parag_table, statistic_chart_title)
    parag_table.add_run(
        f', соотношение по источнику информации {today}',
        style=STYLE
    )
    # parag_table.paragraph_format.right_indent = Inches(0.25)
    parag_table.alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.JUSTIFY

    categories = []
    categories_str = []

    for i in statist_chart_data['smi']:
        if ":" in i['item_date']:
            categories.append(dateutil.parser.parse(i['item_date']).time())
            categories_str.append(f"{dateutil.parser.parse(i['item_date']).time().hour}.00")
        else:
            start_date = dateutil.parser.parse(i['item_date']).date()
            categories.append(start_date)
            categories_str.append(f"{start_date.day}.{start_date.month}")

    # if periods_data.get("period") == "day":
    #     for i in range(today_all.hour + 1):
    #         categories.append(i)
    #         categories_str.append(f"{i}.00")
    # else:
    #     # start_date = get_from_date_datetime(periods_data.get("period")).date()
    #     start_date = dateutil.parser.parse(periods_data.get("_from_data")).date()
    #     while start_date <= dateutil.parser.parse(periods_data.get("_to_data")).date():
    #         categories.append(start_date)
    #         categories_str.append(f"{start_date.day}.{start_date.month}")
    #         start_date += timedelta(days=1)

    negative_list_smi = [0] * len(categories)
    neutral_list_smi = [0] * len(categories)
    positive_list_smi = [0] * len(categories)
    negative_list_social = [0] * len(categories)
    neutral_list_social = [0] * len(categories)
    positive_list_social = [0] * len(categories)

    look_list = [0] * len(categories)

    smi_list = [0] * len(categories)
    social_list = [0] * len(categories)

    for i in range(len(categories)):
        look_list[i] = statist_chart_data['smi'][i]['attendance'] + statist_chart_data['social'][i]['attendance']
        positive_list_smi[i] = statist_chart_data['smi'][i]['positive']
        negative_list_smi[i] = statist_chart_data['smi'][i]['negative']
        neutral_list_smi[i] = statist_chart_data['smi'][i]['netural']
        positive_list_social[i] = statist_chart_data['social'][i]['positive']
        negative_list_social[i] = statist_chart_data['social'][i]['negative']
        neutral_list_social[i] = statist_chart_data['social'][i]['netural']
        smi_list[i] = statist_chart_data['smi'][i]['item_count']
        social_list[i] = statist_chart_data['social'][i]['item_count']

    chart_data = CategoryChartData()
    chart_data.categories = categories_str

    chart_data.add_series('Сми', update_chart_none(smi_list))
    chart_data.add_series('СоцСети', update_chart_none(social_list))
    x, y, cx, cy = Inches(-3.5), Inches(0), Inches(6.15), Inches(3.3)

    chart = document.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
    change_color(chart.plots[0].series[0], RGBColor(255, 134, 13))
    change_color(chart.plots[0].series[1], RGBColor(87, 57, 132))

    update_chart_style(chart)

    if (sum(neutral_list_social) + sum(negative_list_social) + sum(positive_list_social)) > 0:
        if chart_number % 2 == 0:
            document.add_page_break()

        chart_number += 1
        add_table_tonal(document, "в социальных сетях", chart_number, statistic_chart_title, today, categories_str,
                        negative_list_social, neutral_list_social, positive_list_social,
                        x, y, cx, cy)

    # if chart_number % 2 == 1 and period == "day":
    #     parag_table = document.add_paragraph()
    #     parag_table.add_run(
    #         f' ',
    #         style=STYLE
    #     )
    if (sum(negative_list_smi) + sum(neutral_list_smi) + sum(positive_list_smi)) > 0:
        if chart_number % 2 == 0:
            document.add_page_break()
        chart_number += 1

        add_table_tonal(document, "СМИ", chart_number, statistic_chart_title, today, categories_str,
                        negative_list_smi, neutral_list_smi, positive_list_smi,
                        x, y, cx, cy)
    if chart_number % 2 == 0:
        document.add_page_break()
    return chart_number


def add_table_tonal(document, chart_title_type_, chart_number, statistic_chart_title, today, categories_str,
                    negative_list, neutral_list, positive_list,
                    x, y, cx, cy
                    ):
    parag_table = document.add_paragraph()
    parag_table.add_run(
        f' График {chart_number} - Динамика распространения публикаций {chart_title_type_} с упоминанием ',
        style=STYLE
    )
    add_name(parag_table, statistic_chart_title)
    parag_table.add_run(
        f', соотношение по тональности {today}',
        style=STYLE
    )
    parag_table.alignment = docx.enum.text.WD_ALIGN_PARAGRAPH.JUSTIFY

    chart_data = CategoryChartData()
    chart_data.categories = categories_str

    chart_data.add_series('Негативные', negative_list)
    chart_data.add_series('Нейтральные', neutral_list)
    chart_data.add_series('Позитивные', positive_list)
    chart = document.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

    change_color(chart.plots[0].series[0], RGBColor(255, 0, 0))
    change_color(chart.plots[0].series[1], RGBColor(180, 180, 180))
    change_color(chart.plots[0].series[2], RGBColor(0, 255, 0))

    update_chart_style(chart)


if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
